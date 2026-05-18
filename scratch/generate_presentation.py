import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set standard widescreen 16:9 size (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    bg_color = RGBColor(11, 15, 25)        # Deep Blue-Black #0B0F19
    card_color = RGBColor(22, 30, 49)      # Slate Card #161E31
    text_white = RGBColor(255, 255, 255)   # White
    text_grey = RGBColor(148, 163, 184)    # Slate Light Grey #94A3B8
    accent_blue = RGBColor(66, 133, 244)   # Google Blue
    accent_green = RGBColor(52, 168, 83)   # Google Green
    accent_red = RGBColor(234, 67, 53)     # Google Red
    accent_yellow = RGBColor(251, 188, 5)  # Google Yellow

    img_dir = "/Users/kallolchakraborty/Documents/chhanda-local LLM/Presentation"
    
    def set_bg(slide):
        # Draw background rectangle
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="CHHANDA — THE LOCAL AI GATEWAY"):
        # Add category tag
        tag_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.3))
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = accent_blue
        
        # Add main title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(11), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = text_white

    def add_card(slide, left, top, width, height, title="", color=card_color, line_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = color
        if line_color:
            card.line.color.rgb = line_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        
        if title:
            # Shift down slightly to allow title
            title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = text_white
            
        return card

    # ==========================================================
    # SLIDE 1: Title Slide (Cover)
    # ==========================================================
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    
    # Left Content Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CHHANDA (ছন্দা)"
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    p2 = tf.add_paragraph()
    p2.text = "The On-Device Local AI Gateway"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = accent_green
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Harnessing Google's Gemma 4 and LiteRT-LM for 100% Offline, Privacy-First AI Access at Scale."
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = text_grey
    p3.space_before = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = "\n\n• Track: Gemma 4 Good (Global Resilience & Digital Equity)\n• Stack: Kotlin, LiteRT-LM, Ktor-CIO, Room DB, Android KeyStore\n• Author: Solo-Developed by Kallol Chakraborty\n• Dedication: Dedicated to my beloved mother, Chhanda Chakraborty"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.color.rgb = text_grey
    p4.space_before = Pt(20)
    
    # Right Image
    img_path = os.path.join(img_dir, "chhanda_card_thumbnail.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.8), width=Inches(5.0))

    # ==========================================================
    # SLIDE 2: The Core Problem & vision
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "The Rural Digital Divide & On-Device Vision", "THE PROBLEM & THE MISSION")
    
    # Card 1: The Problem
    add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.2), "The Problem: Locked Out of AI")
    prob_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.4), Inches(1.6))
    tf = prob_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• 600 Million Offline: Over half of rural India & Bangladesh lacks consistent internet access, rendering cloud-based AI tools completely useless.\n• Paywalls & Costs: A $20/month ChatGPT subscription is prohibitive in developing economies.\n• Privacy Vulnerability: Uploading patient medical data or private village spreadsheets to overseas cloud engines creates heavy security risks."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey
    
    # Card 2: The Vision
    add_card(slide, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.6), "The Vision: Chhanda & Gemma 4")
    sol_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(5.4), Inches(2.0))
    tf = sol_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• 100% Offline AI: Zero cloud dependence. Gemma 4B is GGUF-quantized and runs natively on the device via LiteRT-LM runtime.\n• One Phone serves 20: Built-in local HTTP server serves an entire classroom or local clinic via Wi-Fi hotspot connectivity.\n• 3-Language Accessibility: Pure native interaction in English, Hindi, and Bengali with customized multilingual Text-To-Speech locale scripting.\n• Zero Cost, Zero Latency: Full edge execution with zero data usage."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey

    # Right Image
    img_path = os.path.join(img_dir, "Screenshot_20260518_041024_Chhanda.jpg")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.8), Inches(1.5), height=Inches(5.0))
        # Add screenshot caption
        cap_box = slide.shapes.add_textbox(Inches(6.8), Inches(6.6), Inches(5.8), Inches(0.4))
        p = cap_box.text_frame.paragraphs[0]
        p.text = "Figure: Multi-client hotspot connectivity configuration guide inside Chhanda Control"
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.color.rgb = accent_green
        p.alignment = PP_ALIGN.CENTER

    # ==========================================================
    # SLIDE 3: Multi-client Ktor-CIO Local Server (LAN Gateway)
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Embedded Ktor-CIO Server: One Phone → 20 Users", "MULTI-CLIENT GATEWAY INFRASTRUCTURE")
    
    # Left Card
    add_card(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.1), "Collaborative Offline Server Node")
    serv_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.2))
    tf = serv_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Local OpenAI Endpoint: Embedded Ktor-CIO HTTP server exposes a standard `/v1/chat/completions` API on the local LAN.\n• Zero-Install Web Client: Students connect to the host hotspot, scan a QR code, and open a beautiful, responsive web portal on any browser to chat with Gemma.\n• IDE & Code Copilot Integration: Software developers in offline areas configure VS Code / Continue to point to port 8888 on the phone for full-fledged code completion.\n• Extreme Concurrency Guard: A native Semaphore system locks concurrent inference to 2 parallel tasks, with a leaky-bucket rate limiter per IP, ensuring a 100% crash-free server."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey
    
    # Right Images (2 side-by-side screenshots)
    img_path1 = os.path.join(img_dir, "Screenshot_20260518_040401_Chhanda.jpg")
    img_path2 = os.path.join(img_dir, "Screenshot_20260518_044400_WebPortal.png")
    
    if os.path.exists(img_path1):
        slide.shapes.add_picture(img_path1, Inches(7.0), Inches(1.5), height=Inches(4.8))
    if os.path.exists(img_path2):
        # Scale to match
        slide.shapes.add_picture(img_path2, Inches(9.5), Inches(2.0), width=Inches(3.3))
        
    cap_box = slide.shapes.add_textbox(Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.4))
    p = cap_box.text_frame.paragraphs[0]
    p.text = "Left: Active local Ktor Server running on device | Right: Responsive Multi-client Web Chat Portal"
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.color.rgb = accent_blue

    # ==========================================================
    # SLIDE 4: On-Device Int8-Quantized RAG Pipeline
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Native Document Ingestion & High-Performance Vector RAG", "KNOWLEDGE RETRIEVAL & VECTOR DATABASE")
    
    add_card(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.1), "100% Private Knowledge Augmentation")
    rag_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.2))
    tf = rag_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Rich Document Support: Zero-cloud offline ingestion of PDF, Word (.docx), Excel (.xlsx), CSV, TSV, HTML, Markdown, JSON, and Images via on-device ML Kit OCR.\n• Int8 Embedding Quantization: Compress high-dimensional Float32 vector embeddings into Int8 space. Reduces database size and memory consumption by 75% (4B/dim → 1B/dim) with negligible accuracy loss.\n• Min-Heap Top-K Search: Implemented native O(N log K) retrieval sorting in Kotlin rather than expensive O(N log N) library arrays, providing 3x faster local search recall.\n• Pronoun-Aware Context Expansion: Dynamically parses follow-up questions containing pronouns and merges them with previous conversation history to guarantee relevant retrieval results."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey
    
    img_path1 = os.path.join(img_dir, "Screenshot_20260518_041133_Chhanda.jpg")
    img_path2 = os.path.join(img_dir, "Screenshot_20260518_041202_Chhanda.jpg")
    
    if os.path.exists(img_path1):
        slide.shapes.add_picture(img_path1, Inches(7.0), Inches(1.5), height=Inches(4.8))
    if os.path.exists(img_path2):
        slide.shapes.add_picture(img_path2, Inches(9.8), Inches(1.5), height=Inches(4.8))
        
    cap_box = slide.shapes.add_textbox(Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.4))
    p = cap_box.text_frame.paragraphs[0]
    p.text = "Left: Offline Ingestion control dashboard | Right: Active Vector storage & RAG recall health trackers"
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.color.rgb = accent_green

    # ==========================================================
    # SLIDE 5: Defense-in-Depth Security Framework
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Defense-In-Depth Security & Android TEE Isolation", "TRUSTED SECURITY ARCHITECTURE")
    
    add_card(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.1), "Enterprise-Grade Threat Defense")
    sec_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.2))
    tf = sec_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• 3-Layer Prompt Injection Shield: A solid filter guarding the local model input context. Blocks system-prompt escape attempts using strict keyword matching, regex rules, and logical structural delimiters (`[USER_INPUT_START/END]`).\n• TEE-Backed Cryptography: Hugging Face read-only download tokens and network access keys are encrypted via Android's hardware keystore in the Trusted Execution Environment (TEE) / Secure Element (SE).\n• PII Redaction & Local Audits: Local filters automatically strip and redact emails, phone numbers, and names prior to passing inputs to the LLM.\n• Secure Biometric Shielding: Complete gate lock protecting server controls behind biometric authentication."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey
    
    img_path1 = os.path.join(img_dir, "Screenshot_20260518_041332_Chhanda.jpg")
    img_path2 = os.path.join(img_dir, "Screenshot_20260518_041045_Chhanda.jpg")
    
    if os.path.exists(img_path1):
        slide.shapes.add_picture(img_path1, Inches(7.0), Inches(1.5), height=Inches(4.8))
    if os.path.exists(img_path2):
        slide.shapes.add_picture(img_path2, Inches(9.8), Inches(1.5), height=Inches(4.8))
        
    cap_box = slide.shapes.add_textbox(Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.4))
    p = cap_box.text_frame.paragraphs[0]
    p.text = "Left: Secure TEE Biometric Authorization shield | Right: API Access Key & node client settings panel"
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.color.rgb = accent_red

    # ==========================================================
    # SLIDE 6: Senior-Grade UI & UX Customization
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Zero-Clutter Visual Design & Premium Customization", "UI/UX & ENGINE MANAGEMENT")
    
    add_card(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.1), "Stunning Material 3 User Experience")
    ux_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.2))
    tf = ux_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Interactive Log Console: Includes a beautiful, live scrolling terminal debugging feed directly in the app, giving developers clear diagnostic insights.\n• Hot-Swappable Local Models: Under 2 seconds! Releasing old engine binders, cleaning JVM garbage collections, and launching new models gracefully without manual app restarts.\n• Multilingual Locale Dynamic TTS: Programmatic real-time script parser. Detects Bengali (`\\u0980-\\u09FF`) and Hindi (`\\u0900-\\u097F`) character blocks and switches active TTS vocal accents automatically.\n• Acoustic Citation Scrubbing: Pre-synthesizer regex parser that strips out bracketed markdown anchors (like `[1]` or `[Source #3]`), verbal lists, and citation text, making speech synthesis incredibly fluid and warm."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey
    
    img_path1 = os.path.join(img_dir, "Screenshot_20260518_041343_Chhanda.jpg")
    img_path2 = os.path.join(img_dir, "Screenshot_20260518_041230_Chhanda.jpg")
    
    if os.path.exists(img_path1):
        slide.shapes.add_picture(img_path1, Inches(7.0), Inches(1.5), height=Inches(4.8))
    if os.path.exists(img_path2):
        slide.shapes.add_picture(img_path2, Inches(9.8), Inches(1.5), height=Inches(4.8))
        
    cap_box = slide.shapes.add_textbox(Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.4))
    p = cap_box.text_frame.paragraphs[0]
    p.text = "Left: Live debug logging terminal dashboard | Right: Theme toggles and TTS Language localizer settings"
    p.font.name = "Arial"
    p.font.size = Pt(9.5)
    p.font.color.rgb = accent_yellow

    # ==========================================================
    # SLIDE 7: Conversational Live Chat & Performance
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Real-Time Streaming Chat & Performance Telemetry", "CONVERSATIONAL EXPERIENCE & PERFORMANCE")
    
    add_card(slide, Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.1), "Fluent Interface & Telemetry Widgets")
    chat_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.1), Inches(4.2))
    tf = chat_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Micro-Indicator Feedback: Beautiful chat bubbles showing dynamic tokens-per-second, generation count, and exact sub-millisecond RAG query latency directly underneath message bubbles.\n• Hands-Free Single-Shot Voice: Speaks responses clearly and resets gracefully, maintaining a 500ms safety lock to protect against audio-loop echo triggers.\n• Thermal-Aware Dynamic Contexts: Listens to Android hardware thermal status trackers. Reduces the active model context window size dynamically to protect the phone against hot throttling or high JNI memory crashes.\n• Lifecycle-Aware Telemetry: Shuts down background polling monitors when the app is minimized (`onStop`), saving 12% idle battery drainage."
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.color.rgb = text_grey
    
    img_path = os.path.join(img_dir, "Screenshot_20260518_041119_Chhanda.jpg")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8.0), Inches(1.5), height=Inches(4.8))
        
    cap_box = slide.shapes.add_textbox(Inches(8.0), Inches(6.5), Inches(4.5), Inches(0.4))
    p = cap_box.text_frame.paragraphs[0]
    p.text = "Figure: Streaming Chat UI with micro-performance stats"
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.color.rgb = accent_blue
    p.alignment = PP_ALIGN.CENTER

    # ==========================================================
    # SLIDE 8: Competitive Edge Matrix
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Competitive Advantage Analysis", "MARKET ALIGNMENT & HARDENED METRICS")
    
    # Draw table
    rows, cols = 9, 4
    left, top, width, height = Inches(0.6), Inches(1.6), Inches(12.13), Inches(5.0)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(3.53) # Feature
    table.columns[1].width = Inches(2.8)  # Google Edge Gallery
    table.columns[2].width = Inches(2.8)  # ChatGPT Mobile
    table.columns[3].width = Inches(3.0)  # Chhanda AI Gateway
    
    # Table data
    headers = ["FEATURE / METRIC", "GOOGLE AI EDGE GALLERY", "CHATGPT (MOBILE)", "CHHANDA AI GATEWAY"]
    data = [
        ["Fully Offline Operation", "Partial (No server capability)", "Requires Cloud Connection", "✅ 100% Offline (Local Gemma 4)"],
        ["On-Device Int8 Vector RAG", "❌ No Vector Storage", "❌ Cloud RAG only", "✅ Yes (PDF/Spreadsheets/OCR)"],
        ["Multi-Client Local Server", "❌ Single User Only", "❌ Single Account Only", "✅ Yes (Ktor-CIO, up to 20 users)"],
        ["OpenAI API Compatibility", "❌ Proprietary bindings", "✅ Cloud Only", "✅ Yes (Local REST Gateway)"],
        ["Prompt Injection Defense", "❌ No Local Guards", "✅ Cloud Guardrails Only", "✅ Yes (3-Layer Prompt Shield)"],
        ["Thermal Auto-Throttle", "❌ None", "N/A", "✅ Yes (Dynamic Context Scale)"],
        ["Localized Multilingual Speech", "❌ English Only", "✅ Cloud synthesis", "✅ Yes (EN/HI/BN with Dynamic Script)"],
        ["Commercial Software Cost", "Free (Dev sample)", "$20/month Plus subscription", "✅ 100% Free Forever"],
    ]
    
    # Style Header
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = card_color
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = accent_blue
        p.alignment = PP_ALIGN.CENTER
        
    # Fill Data
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 23, 42) if row_idx % 2 == 0 else RGBColor(22, 30, 49)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.name = "Arial"
            p.font.size = Pt(10)
            p.font.color.rgb = text_white
            if col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = text_white
            elif col_idx == 3:
                p.font.color.rgb = accent_green
                p.font.bold = True
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    # ==========================================================
    # SLIDE 9: "Gemma 4 Good" Social Impact
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Democratizing State-of-the-Art Offline Intelligence", "GEMMA 4 GOOD SOCIAL IMPACT")
    
    # 3 horizontal cards
    card_w = Inches(3.7)
    card_h = Inches(4.5)
    gap = Inches(0.4)
    y_pos = Inches(1.8)
    
    # Card 1: Education
    add_card(slide, Inches(0.6), y_pos, card_w, card_h, "🌍 Inclusive Education")
    ed_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), card_w - Inches(0.4), Inches(3.6))
    tf = ed_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Offline Infrastructure: A single ₹15,000 Android phone running Chhanda serves as the local AI library for a rural school.\n• Interactive Ingestion: Teachers upload textbook sets and mock assessments into local vector databases for direct offline student Q&A.\n• Localized Learning: Multilingual interaction in regional languages breaks major literacy barriers."
    p.font.name = "Arial"
    p.font.size = Pt(11.5)
    p.font.color.rgb = text_grey
    
    # Card 2: Healthcare
    add_card(slide, Inches(0.6) + card_w + gap, y_pos, card_w, card_h, "🏥 Rural Healthcare")
    hc_box = slide.shapes.add_textbox(Inches(0.8) + card_w + gap, Inches(2.4), card_w - Inches(0.4), Inches(3.6))
    tf = hc_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Absolute Patient Privacy: Zero data leakage. Healthcare workers analyze sensitive medical symptoms and records entirely offline.\n• Fast Reference Lookup: Nurses query heavy clinical guidelines instantly via RAG vector indices under network dropouts.\n• On-Device Safety Shields: PII auto-redactors scrub identifying elements automatically."
    p.font.name = "Arial"
    p.font.size = Pt(11.5)
    p.font.color.rgb = text_grey
    
    # Card 3: Digital Equity
    add_card(slide, Inches(0.6) + (card_w + gap)*2, y_pos, card_w, card_h, "🤝 Complete Digital Equity")
    eq_box = slide.shapes.add_textbox(Inches(0.8) + (card_w + gap)*2, Inches(2.4), card_w - Inches(0.4), Inches(3.6))
    tf = eq_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "• Zero Cost Utility: Replaces heavy monthly commercial APIs, making AI open and accessible to the last mile.\n• Privacy Sovereignty: Promotes 'Your AI, Your Data, Your Device' in an era of mass data capture.\n• Mobile-First Resilience: Thrives in zero-connectivity disaster zones, blackouts, and remote regions."
    p.font.name = "Arial"
    p.font.size = Pt(11.5)
    p.font.color.rgb = text_grey

    # ==========================================================
    # SLIDE 10: Conclusion & Call to Action
    # ==========================================================
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "The Future of Edge AI is Private and Offline", "CONCLUSION")
    
    # Large Center Card
    add_card(slide, Inches(0.6), Inches(1.6), Inches(12.13), Inches(4.8), "Chhanda — Redefining Global Accessibility", color=card_color, line_color=accent_green)
    
    conc_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.8))
    tf = conc_box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    
    p = tf.paragraphs[0]
    p.text = "• Complete Technological Sovereignty: 14,500+ lines of Kotlin delivering complete independence from international cloud paywalls."
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = text_white
    
    p2 = tf.add_paragraph()
    p2.text = "• Powered by Gemma 4: Deep native execution utilizing Google AI Edge's LiteRT-LM libraries directly on mobile hardware."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = text_white
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "• Absolute Last-Mile Inclusivity: One single smartphone hosting a local Wi-Fi gateway servers up to 20 concurrent browsers offline."
    p3.font.name = "Arial"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = text_white
    p3.space_before = Pt(12)
    
    p4 = tf.add_paragraph()
    p4.text = "\n\"Chhanda represent the beautiful harmony between the structural limits of on-device code and the endless, flowing possibilities of localized human intelligence.\""
    p4.font.name = "Arial"
    p4.font.size = Pt(14)
    p4.font.italic = True
    p4.font.color.rgb = text_grey
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(25)
    
    p5 = tf.add_paragraph()
    p5.text = "\nThank You! | Solo Developed by Kallol Chakraborty | GitHub: https://github.com/kallolchakraborty/Chhanda-AI"
    p5.font.name = "Arial"
    p5.font.size = Pt(12)
    p5.font.bold = True
    p5.font.color.rgb = accent_blue
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(10)

    # Save
    out_path = os.path.join(img_dir, "Chhanda_Gemma_4_Good_Presentation.pptx")
    prs.save(out_path)
    print(f"Presentation saved successfully to: {out_path}")

if __name__ == "__main__":
    create_presentation()
